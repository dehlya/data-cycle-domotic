"""
polish_pptx_dehlya.py — improve slides 7-13 (Bronze to Gold) of the team
deck without touching any other slide.

Reads:  C:\\Users\\dehly\\Downloads\\Data-Cycle-Project 1.pptx
Writes: docs\\v2\\out\\Data-Cycle-Project_polished.pptx

Strategy: open with python-pptx, iterate text frames on slides 7-13 only,
apply targeted run-level text replacements (preserves all formatting,
fonts, colours, positions, images, and the visual template).
"""

from copy import deepcopy
from pathlib import Path
from pptx import Presentation

INPUT_PATH  = Path(r"C:\Users\dehly\Downloads\Data-Cycle-Project 1.pptx")
OUTPUT_PATH = Path("docs/v2/out/Data-Cycle-Project_polished.pptx")
TARGET_SLIDES = set(range(7, 14))  # 7, 8, 9, 10, 11, 12, 13


# (substring_to_find, replacement_text)
# Applied to ANY text frame on slides 7-13. Keep changes specific so we
# don't accidentally touch unintended copy.
REPLACEMENTS = [
    # ── SLIDE 7: Bronze ─────────────────────────────────────────────────────
    # Opening paragraph: solid; tighten the closing tagline
    ("Bronze is the traceability and safety layer of the project.",
     "Bronze is the traceability and safety net of the pipeline — every other layer can be rebuilt from it."),

    # ── SLIDE 8: Bronze → Silver ────────────────────────────────────────────
    # Original opening had repeating "Silver transforms..." phrasing and was
    # truncated. Replace with a tighter version.
    ("The Silver layer transforms raw data into clean, typed and structured "
     "relational data. Silver transforms raw technical files into reliable "
     "structured data. Every later output depends on the quality of",
     "The Silver layer turns raw files into clean, typed, structured "
     "relational data. Every later output — Gold tables, Power BI "
     "dashboards, ML predictions — depends on the quality of this "
     "transformation."),

    # Second-level fix in case the truncation differs in the file
    ("Silver transforms raw technical files into reliable structured data.",
     ""),

    # ── SLIDE 9: Silver Data Model ──────────────────────────────────────────
    ("The Silver layer contains normalised operational data. It is designed "
     "to represent the cleaned version of the source data. Silver is the "
     "clean operational foundation of the project.",
     "The Silver layer is the clean operational foundation of the project — "
     "normalised, typed tables that mirror the source systems with errors "
     "and inconsistencies removed."),

    ("The Silver layer is not yet the final analytical model, but it provides "
     "the reliable base required to build that model.",
     "Silver is not yet the analytical model — it is the reliable base that "
     "Gold is built on top of."),

    # ── SLIDE 10: Silver → Gold ─────────────────────────────────────────────
    ("The Gold layer transforms cleaned data into a business-oriented "
     "analytical model. Gold is the bridge between cleaned data and final "
     "exploitation — designed for BI and ML, not only storage.",
     "The Gold layer turns cleaned data into a business-oriented analytical "
     "model. Designed for BI and ML — not just storage — it is the bridge "
     "between technical data and final exploitation."),

    ("The project uses a star schema composed of two types of tables:",
     "The Gold model uses a star schema with two table families:"),

    # ── SLIDE 11: Gold Dimensions ───────────────────────────────────────────
    ("Dimensions provide the descriptive context required for analysis. "
     "Without dimensions, facts are only numbers. Dimensions allow users to "
     "understand where, when and by which device each measurement was",
     "Dimensions are the descriptive context for every measurement. "
     "Without them, facts are only numbers — dimensions tell you where, "
     "when, and by which device each measurement was captured."),

    # Improve dim descriptions — be more specific
    ("Provides date and time attributes for time-based analysis. Enables "
     "filtering by hour, day, week, month or year.",
     "Date and time attributes at minute granularity. Enables filtering by "
     "hour, day, week, month, weekday or year."),

    ("Describes each apartment. Allows dashboards to compare behaviour across "
     "different apartment units.",
     "Apartment metadata (anonymised). Lets dashboards compare behaviour "
     "across apartments while preserving Row-Level Security."),

    ("Describes rooms and their relationship with apartments. Enables "
     "room-level granularity in all analyses.",
     "Rooms with their parent apartment. Enables room-level granularity "
     "across every fact table."),

    ("Describes sensors and devices. Allows filtering and grouping by device "
     "type, location and status.",
     "Sensors and devices linked to their room. Used to filter and group by "
     "device type, location, and operational status."),

    ("Dimensions allow dashboards to filter, group and compare data by time, "
     "apartment, room or device — giving business meaning to numerical "
     "measurements.",
     "Together, dimensions give business meaning to numerical measurements — "
     "the same fact can be sliced by time, apartment, room, or device."),

    # ── SLIDE 12: Gold Fact Tables ──────────────────────────────────────────
    ("Power and energy consumption indicators. Core metrics for monitoring "
     "apartment electricity usage.",
     "Power (W) and energy (kWh) per device per minute. Cost projections via "
     "the tariff dimension."),

    ("Temperature, humidity, CO₂ and other comfort indicators. Tracks indoor "
     "environmental quality over time.",
     "Temperature, humidity, CO₂, noise and pressure per room per minute. "
     "Outliers flagged at the silver layer."),

    ("Room occupancy and motion-based presence indicators. Captures when and "
     "where people are detected.",
     "Motion, door and window flags per room per minute. Foundation for the "
     "occupancy-prediction ML workflow."),

    ("Sensor activity, reliability and data availability indicators. "
     "Monitors the health of the IoT infrastructure.",
     "Daily uptime, error counts, missing readings and battery levels per "
     "device. Surfaces failing sensors before they impact the analysis."),

    ("Machine learning prediction outputs. Stores model results alongside "
     "actuals for comparison and evaluation.",
     "Motion and consumption forecasts written by KNIME, with the prediction "
     "timestamp recorded so model versions can be compared."),

    # ── SLIDE 13: Use Cases ─────────────────────────────────────────────────
    # Punctuation: existing items end with ; and last with . — keep as is,
    # but fix any inconsistencies and lower-case awkwardness
    ("avg power consumption ;",            "Average power consumption ;"),
    ("peak consumption ;",                 "Peak consumption ;"),
    ("consumption by room.",               "Consumption by room."),
    ("missing data ;",                     "Missing data ;"),
    ("inactive devices ;",                 "Inactive devices ;"),
    ("last seen timestamp.",               "Last-seen timestamp."),
    ("occupied time by room ;",            "Occupied time by room ;"),
    ("presence by hour ;",                 "Presence by hour ;"),
    ("weekday/weekend patterns.",          "Weekday vs weekend patterns."),
    ("predicted occupancy ;",              "Predicted occupancy ;"),
    ("predicted consumption ;",            "Predicted consumption ;"),
    ("prediction vs actual.",              "Predicted vs actual."),
    ("average temperature ;",              "Average temperature ;"),
    ("humidity range ;",                   "Humidity range ;"),
    ("CO₂ level.",                         "CO₂ level."),  # already correct, keep

    # Slide 13 lead-in — slightly tighten
    ("The Gold layer supports both descriptive analytics and predictive "
     "analytics. This slide connects the data model with the final business "
     "value of the project.",
     "The Gold layer supports both descriptive and predictive analytics — "
     "this is where the data model meets real business value."),

    # ── Cleanup pass: trailing fragments left from multi-run replacements ──
    # Slide 8: stray "S" left behind from "Silver transforms..." sentence
    ("transformation. Silver", "transformation."),
    ("transformation. S",       "transformation."),
    ("transformation. data.",   "transformation."),
    ("transformation.data.",    "transformation."),
    ("transformation. Data.",   "transformation."),
    # Slide 11: dangling "produced." left after replacing the truncated
    # "...measurement was [truncated] produced." sentence
    ("captured. produced.", "captured."),
    ("captured.produced.",  "captured."),

    # ── Bronze compression detail (slide 7 closing) ─────────────────────────
    ("Bronze is the traceability and safety net of the pipeline — every other layer can be rebuilt from it.",
     "After silver ingests a file, bronze gzips it in place — disk drops ~10-15× while the audit trail stays intact. Every other layer can be rebuilt from bronze."),

    # ── Performance note on Silver hop (slide 8 — adds context to step 04) ──
    ("Insert clean records into relational database tables ready for querying and joining.",
     "Bulk-load via COPY into a temp table + single INSERT ... ON CONFLICT (50-150× faster than per-row INSERT)."),

    # ── Slide 9: Time References doesn't exist in silver. Fix to a real entity ──
    # silver.di_errors_clean is populated from MySQL DIErrors and joined into
    # gold.fact_device_health_day. Calling it "Time References" was misleading.
    ("Time References", "Sensor Error Logs"),
]


# ── Speaker notes for slides 7-13 ────────────────────────────────────────────
# ~50 seconds of talking per slide ≈ 100-130 words. Covers the slide content
# + adds context the audience can't read on the slide.
SPEAKER_NOTES = {
    # ── Slide 7 — BRONZE ─────────────────────────────────────────────────────
    7: (
        "OPEN — \"The hardest part of an IoT pipeline isn't the dashboards "
        "— it's keeping the raw data trustworthy from the moment it lands.\"\n\n"
        "Two apartments × one JSON per minute = 400 735 files in our "
        "current backfill. Plus one daily weather CSV on the sFTP.\n\n"
        "ENGINEERING DETAILS:\n"
        "• Path layout: storage/bronze/jimmy/2023/10/15/14/<file>.json — "
        "year/month/day/hour folders so even after years no folder has "
        "more than ~60 files. NTFS handles that fast; flat trees would "
        "choke.\n"
        "• bulk_to_bronze.py runs in PREDICTIVE mode by default: from "
        "the newest filename in bronze it computes the next expected one, "
        ".exists() on the SMB — about 5 ms per check. Falls back to a "
        "full scan once a night or when bronze is empty.\n"
        "• After silver successfully ingests a file, the watcher gzips "
        "the bronze copy in place — typical sensor JSON shrinks 10-15×, "
        "audit trail intact. Filename also appended to processed.log so "
        "future SMB rescans skip it without a Postgres roundtrip.\n\n"
        "TRANSITION — \"That's the safety net. Now the heavy lifting.\""
    ),
    # ── Slide 8 — BRONZE → SILVER ───────────────────────────────────────────
    8: (
        "OPEN — \"Silver is the heaviest hop and the place where I spent "
        "most of the engineering effort.\"\n\n"
        "ENGINEERING DETAILS:\n"
        "• flatten_sensors.py uses ProcessPoolExecutor with 8 workers, "
        "batches of 2 000 files. Each worker parses JSON, normalises room "
        "names (\"Bdroom\" → \"Bedroom\"), applies outlier bounds (e.g. "
        "temperature ∈ [-20, 60]°C, CO₂ ∈ [300, 5000] ppm), then returns "
        "rows to the main process.\n"
        "• THE TRICK: instead of per-row INSERT, we use psycopg2's "
        "copy_expert to stream rows into a TEMP TABLE … ON COMMIT DROP, "
        "then merge with one statement:\n"
        "      INSERT INTO silver.sensor_events …\n"
        "      SELECT DISTINCT ON (apartment, room, sensor_type, field, timestamp) …\n"
        "      ON CONFLICT (…) DO UPDATE SET …\n"
        "  DISTINCT ON dedupes within the batch so Postgres doesn't "
        "error on the same key appearing twice. One round-trip, "
        "set-based, sustained 10-30 k rows/sec.\n"
        "• Idempotency comes from silver.etl_watermark — every filename "
        "we ingest gets inserted there. find_new_files_fast() diffs "
        "against this set; re-runs only process files NOT in the "
        "watermark. No double-counting, ever.\n\n"
        "RESULT — original backfill: 3 hours. After this change: "
        "10-15 minutes (50-150× speedup, measured). Without it, the "
        "for-dummies installer wouldn't be viable.\n\n"
        "TRANSITION — \"Reliable structured data. Let's see what's in "
        "silver.\""
    ),
    # ── Slide 9 — SILVER DATA MODEL ─────────────────────────────────────────
    9: (
        "All in PostgreSQL under the silver schema.\n\n"
        "ENGINEERING DETAILS:\n"
        "• From MySQL: silver.dim_buildings, dim_rooms, dim_devices, "
        "dim_sensors. ~16 sensors, 11 rooms, 9 devices across 2 apartments.\n"
        "• From JSON files: silver.sensor_events in LONG format — one "
        "row per reading: (apartment, room, sensor_type, field, value, "
        "unit, timestamp, is_outlier). Why long instead of pivoted? "
        "Sensors come and go; columns shouldn't. Adding a new sensor "
        "type means new ROWS, not a schema migration. ~10 M rows after "
        "a full backfill.\n"
        "• From sFTP: silver.weather_forecasts — ~150 k rows per daily "
        "CSV, kept flat (no pivot) for the same reason.\n"
        "• From MySQL DIErrors: silver.di_errors_clean — typed, with "
        "apartment mapping (idBuilding → \"jimmy\" / \"jeremie\") and "
        "a severity heuristic (regex on the message: fatal/crash/down "
        "→ \"high\", fail/error → \"medium\", else \"low\"). 43 103 "
        "rows in our current install.\n"
        "• Two watermark tables: silver.etl_watermark, "
        "silver.weather_watermark. Every successfully ingested filename "
        "lands there. Re-runs diff against these sets → fully idempotent.\n\n"
        "TRANSITION — \"From operational to analytical: star schema.\""
    ),
    # ── Slide 10 — SILVER → GOLD ────────────────────────────────────────────
    10: (
        "Gold is where silver becomes business-ready.\n\n"
        "ENGINEERING DETAILS:\n"
        "• populate_gold.py is the orchestrator: with --sensors or "
        "--weather flags it triggers populate_dimensions → "
        "populate_sensors → populate_weather in order. Every step is "
        "ONE set-based SQL statement, not a row-by-row loop:\n"
        "      INSERT INTO gold.fact_energy_minute …\n"
        "      SELECT … FROM silver.sensor_events se\n"
        "      JOIN gold.dim_apartment a ON …\n"
        "      JOIN gold.dim_room r ON …\n"
        "      GROUP BY … ON CONFLICT (datetime_key, device_key) "
        "DO UPDATE\n"
        "  Postgres aggregates millions of rows server-side, atomically.\n"
        "• Refresh cadence is built into the watcher: populate_gold "
        "--sensors fires every 15 min, --weather and KNIME predictions "
        "run daily at 06:30. Dashboards stay within 15 min of fresh "
        "without a separate scheduler.\n"
        "• mv_energy_with_cost is a materialised view joining "
        "fact_energy_minute with dim_tariff. REFRESH MATERIALIZED VIEW "
        "CONCURRENTLY at the end of each gold pass, with a non-concurrent "
        "fallback for first-time empty MVs.\n\n"
        "ONE-LINE SUMMARY — \"Silver answers WHAT happened. Gold "
        "answers SO WHAT.\""
    ),
    # ── Slide 11 — GOLD DIMENSIONS ──────────────────────────────────────────
    11: (
        "Four dimensions on the slide; two more behind the scenes. "
        "All with INTEGER surrogate keys + B-tree indexes for fast joins.\n\n"
        "ENGINEERING DETAILS:\n"
        "• dim_datetime — every minute since 2023 (~2.6 M rows). "
        "datetime_key is encoded YYYYMMDDHHMM as a single BIGINT — same "
        "pattern as date_key (YYYYMMDD as INT). Hour, weekday, "
        "is_business_hour pre-computed → Power BI filters \"weekday "
        "evenings\" without a date function.\n"
        "• dim_apartment has anonymisation in the populate_dimensions.py "
        "SQL itself: \n"
        "      WHERE LOWER(b.\"houseName\") LIKE '%jimmy%' → 'Building 1'\n"
        "      WHERE LOWER(b.\"houseName\") LIKE '%jeremie%' → 'Building 2'\n"
        "      owner_user_id is NULLed out\n"
        "  We keep only first names — GDPR Article 4(1): a first name "
        "alone isn't personal data without additional context.\n"
        "• dim_room → parent apartment_key. dim_device → device_id "
        "synthesised as apartment || '_' || room || '_' || sensor_type "
        "(stable natural key from JSON content).\n"
        "• Behind the scenes: dim_tariff (electricity rate per "
        "hour-of-day, drives mv_energy_with_cost), dim_weather_site "
        "(joins outdoor weather to the consumption-prediction model)."
    ),
    # ── Slide 12 — GOLD FACT TABLES ─────────────────────────────────────────
    12: (
        "Five fact tables, one per analytical domain. Each has a UNIQUE "
        "constraint on (datetime_key, device_key) or (datetime_key, "
        "room_key) — guarantees idempotency on re-runs.\n\n"
        "ENGINEERING DETAILS per fact:\n"
        "• fact_energy_minute (W + kWh per device per minute). Built "
        "with MAX(CASE WHEN field='power' THEN value END) + MAX(CASE "
        "WHEN field='total' THEN value/1000 END) — pivots the long "
        "silver into wide gold in one SQL statement. Joined with "
        "dim_tariff in mv_energy_with_cost for CHF.\n"
        "• fact_environment_minute (indoor T, CO₂, noise, pressure per "
        "room per minute). is_anomaly flag rolls up the sensor-level "
        "is_outlier from silver.\n"
        "• fact_presence_minute (motion, door, window flags). presence_flag "
        "uses BOOL_OR over motion-type sensors only — not door, because "
        "door sensors heartbeat too often and would falsely report "
        "occupancy. Real bug we fixed during dev.\n"
        "• fact_device_health_day. CTE-based: battery from "
        "silver.sensor_events, error_count from silver.di_errors_clean, "
        "missing-readings as 1440 minus distinct minutes seen that day.\n"
        "• fact_prediction_motion + fact_prediction_consumption. KNIME "
        "writes directly via DB Writer. prediction_made_at column lets us "
        "compare model versions over time — 66 186 consumption rows + "
        "13 173 motion rows verified in the current install."
    ),
    # ── Slide 13 — USE CASES ────────────────────────────────────────────────
    13: (
        "OPEN — \"This is where the data model becomes business value.\"\n\n"
        "ENGINEERING DETAILS:\n"
        "• Each KPI on this slide is ONE query — a fact table joined "
        "with 2-3 dimensions on integer surrogate keys. Sub-second "
        "in Power BI.\n"
        "• Per-apartment isolation enforced via Row Level Security "
        "in the .pbix model. RLS rule: \"[apartment_key] = X\". "
        "A tenant viewing as Jimmy literally cannot see Jeremie's data — "
        "Power BI rewrites every query with the WHERE clause before "
        "sending it to Postgres. Not just visual hiding.\n"
        "• Refresh: gold updates every 15 min via the watcher; KNIME "
        "predictions land daily at 06:30. Power BI Desktop loads a "
        "snapshot at refresh time — Ctrl+R pulls latest. Documented in "
        "the user guide.\n"
        "• Predictions vs actuals on the right: KNIME writes "
        "fact_prediction_motion / fact_prediction_consumption with a "
        "prediction_made_at timestamp; a single SQL query joins those "
        "rows back to the same minute in fact_energy_minute or "
        "fact_presence_minute → side-by-side comparison, model drift "
        "tracked over time.\n\n"
        "WRAP-UP — \"Three layers, six dimensions, seven fact tables, "
        "two ML models. Every join is on an integer surrogate key. "
        "That's the data engineering foundation.\"\n\n"
        "HANDOFF — \"Sacha takes it from here.\""
    ),
}


def replace_in_text_frame(tf, old, new):
    """Replace 'old' with 'new' inside a text frame.
    Tries run-level first to preserve formatting; falls back to paragraph
    rewrite if 'old' spans multiple runs."""
    n_replaced = 0
    for paragraph in tf.paragraphs:
        # Fast path: substring fits in a single run
        for run in paragraph.runs:
            if old in (run.text or ""):
                run.text = run.text.replace(old, new)
                n_replaced += 1

        # Slow path: substring may span runs
        joined = "".join(r.text or "" for r in paragraph.runs)
        if old in joined and n_replaced == 0:
            new_joined = joined.replace(old, new)
            if paragraph.runs:
                paragraph.runs[0].text = new_joined
                for r in paragraph.runs[1:]:
                    r.text = ""
            n_replaced += 1
    return n_replaced


def main():
    if not INPUT_PATH.exists():
        raise SystemExit(f"Input not found: {INPUT_PATH}")
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

    print(f"Reading  {INPUT_PATH}")
    p = Presentation(INPUT_PATH)

    print(f"\nTarget slides: {sorted(TARGET_SLIDES)}\n")
    total_replacements = 0

    for slide_idx, slide in enumerate(p.slides, start=1):
        if slide_idx not in TARGET_SLIDES:
            continue

        slide_replacements = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for old, new in REPLACEMENTS:
                slide_replacements += replace_in_text_frame(shape.text_frame, old, new)

        if slide_replacements:
            print(f"  Slide {slide_idx:2d}: {slide_replacements} replacement(s)")
        total_replacements += slide_replacements

        # Set/replace speaker notes for this slide (notes pane in PowerPoint)
        notes_text = SPEAKER_NOTES.get(slide_idx)
        if notes_text:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.clear()
            tf.text = notes_text

    print(f"\n{total_replacements} text replacement(s) applied to slides 7-13.")
    print(f"{len(SPEAKER_NOTES)} speaker-notes blocks written.")
    p.save(OUTPUT_PATH)
    size_kb = OUTPUT_PATH.stat().st_size / 1024
    print(f"\n✓ Wrote {OUTPUT_PATH} ({size_kb:.0f} KB)")
    print(f"  Slides 1-6 and 14-21 untouched.")


if __name__ == "__main__":
    main()
